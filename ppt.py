from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation():
    prs = Presentation()
    
    # --- Slide 1: Title Slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide1 = prs.slides.add_slide(title_slide_layout)
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Aletheia - Explainable Image Forensics"
    subtitle.text = "Decoding Synthetic Media through Statistical Feature Extraction\nPresented by: Misky"

    # --- Slide 2: System Architecture & Dataset ---
    bullet_slide_layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    shapes2 = slide2.shapes
    shapes2.title.text = "System Architecture & Dataset"
    body_shape2 = shapes2.placeholders[1]
    tf2 = body_shape2.text_frame
    
    tf2.text = "Kaggle Dataset: ~1500 images (Unsplash vs. Stable Diffusion)"
    p = tf2.add_paragraph()
    p.text = "5 Categories: Nature, City, Animals, Food, People"
    p.level = 1
    p2 = tf2.add_paragraph()
    p2.text = "Full-Stack Architecture:"
    p3 = tf2.add_paragraph()
    p3.text = "Frontend: React UI for real-time user uploads"
    p3.level = 1
    p4 = tf2.add_paragraph()
    p4.text = "Backend: Flask API for model serving"
    p4.level = 1
    p5 = tf2.add_paragraph()
    p5.text = "Core: Python Machine Learning Pipeline"
    p5.level = 1

    # --- Slide 3: Feature Engineering ---
    slide3 = prs.slides.add_slide(bullet_slide_layout)
    shapes3 = slide3.shapes
    shapes3.title.text = "Feature Engineering (The Core Engine)"
    body_shape3 = shapes3.placeholders[1]
    tf3 = body_shape3.text_frame
    
    tf3.text = "Converting raw pixels into statistical data to detect synthetic artifacts:"
    p = tf3.add_paragraph()
    p.text = "Local Binary Patterns (LBP): Analyzes texture smoothness."
    p = tf3.add_paragraph()
    p.text = "Fast Fourier Transform (FFT): Detects high-frequency noise discrepancies."
    p = tf3.add_paragraph()
    p.text = "Gray Level Co-occurrence Matrix (GLCM): Measures contrast and homogeneity."

    # --- Slide 4: Model Benchmarking & XAI ---
    slide4 = prs.slides.add_slide(bullet_slide_layout)
    shapes4 = slide4.shapes
    shapes4.title.text = "Model Benchmarking & Explainable AI"
    body_shape4 = shapes4.placeholders[1]
    tf4 = body_shape4.text_frame
    
    tf4.text = "Evaluated Models: Logistic Regression, Decision Trees, Random Forest"
    p = tf4.add_paragraph()
    p.text = "Focus on Explainable AI (XAI) using Decision Trees."
    p = tf4.add_paragraph()
    p.text = "Moves beyond 'black box' predictions to human-readable rules."
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Example: 'If high-frequency noise is low AND texture variance is smooth -> AI-Generated'."
    p.level = 1

    # --- Slide 5: Dashboard & Key Takeaways ---
    slide5 = prs.slides.add_slide(bullet_slide_layout)
    shapes5 = slide5.shapes
    shapes5.title.text = "Live Dashboard & Key Takeaways"
    body_shape5 = shapes5.placeholders[1]
    tf5 = body_shape5.text_frame
    
    tf5.text = "User Interface allows instant authenticity verification."
    p = tf5.add_paragraph()
    p.text = "Engineering Challenges Overcome:"
    p = tf5.add_paragraph()
    p.text = "Mitigated 'Shortcut Learning' where the model classified camera noise instead of actual features."
    p.level = 1
    p = tf5.add_paragraph()
    p.text = "Ensured robust evaluation using Stratified K-Fold Cross-Validation."
    p.level = 1

    # Save the file
    prs.save('Aletheia_Presentation.pptx')
    print("Success! 'Aletheia_Presentation.pptx' has been generated in your folder.")

if __name__ == '__main__':
    create_presentation()